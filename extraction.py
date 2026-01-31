from pptx import Presentation
import pandas as pd
import json
import re

# ファイルパスの設定
pptx_file = "かずこの会_相関図（ラジオ）.pptx"
output_file = "data/members_radio.json"

# Debugging: Print the PPTX file being processed
print(f"Processing PPTX file: {pptx_file}")

# プレゼンテーションの読み込み
presentation = Presentation(pptx_file)

# Debugging: Confirm the PPTX file and slide count
print(f"Processing PPTX file: {pptx_file}")
print(f"Number of slides: {len(presentation.slides)}")

# Ensure data list is initialized at the start of the script
data = []

# Add a mapping for full location names
location_mapping = {
    "神": "神戸",
    "大": "大阪",
    "滋": "滋賀",
    "兵": "兵庫",
    "奈": "奈良",
    "京": "京都",
    "埼": "埼玉",
    "千": "千葉",
    "愛": "愛媛",
    "福": "福島",
    "秋": "秋田",
    "淡": "淡路島",
    "北": "北海道",
    "パ": "パリ",
    "イ": "イギリス",
    "ミ": "ミナミ",
    "西": "西脇",
    "上": "上本町"
}

# スライド内のテーブルを解析
for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            for row in table.rows:
                try:
                    cells = [cell.text.strip().replace("\n", "") for cell in row.cells]

                    if len(cells) >= 2:
                        name_age_match = re.match(r"(.*?)(\d+歳|\d+代)?$", cells[1])
                        name = name_age_match.group(1).strip() if name_age_match else cells[1]
                        age = name_age_match.group(2) if name_age_match and name_age_match.group(2) else ""

                        comment = ""
                        if len(cells) > 3:
                            comment = cells[3]

                        # Extract the location key and map to full name
                        location_key = cells[0][4:5]  # Extract the 5th character as the location key
                        full_location = location_mapping.get(location_key, location_key)

                        data.append({
                            "会員番号": cells[0][:4],  # Remove the appended character
                            "居住地": full_location,  # Use the full location name
                            "名前": name,
                            "年齢": age,
                            "コメント": comment
                        })
                    else:
                        print("Skipped row due to insufficient columns.")
                except Exception as e:
                    print(f"Error processing row: {e}")

        # Extract comments from text frames
        elif shape.has_text_frame:
            text = shape.text.strip()
            if len(text) > 5:  # Consider as a comment if text is long enough
                if data:
                    data[-1]["コメント"] = text

# デバッグ用: 抽出したデータを出力
print("--- 抽出したデータ ---")
print(data)

# Debugging: Print data and sorted_data before saving
print("--- Final data before saving ---")
print(data)
print("--- Final sorted_data before saving ---")
sorted_data = sorted(data, key=lambda x: x["会員番号"])

# JSONファイルに保存
with open(output_file, "w", encoding="utf-8") as f:
    json.dump(data, f, ensure_ascii=False, indent=4)

print(f"抽出したデータを {output_file} に保存しました。")